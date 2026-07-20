"""Generate the MindCare project presentation (12 slides) as a .pptx file.

Run:  python generate_slides.py
Output:  MindCare_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Theme ---
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT = RGBColor(0xF3, 0xF0, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x0F, 0x9D, 0x58)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return tb


def content_slide(title, kicker, bullets):
    """bullets: list of (text, level) where level 0 = main, 1 = sub."""
    slide = prs.slides.add_slide(BLANK)
    # left accent bar
    add_rect(slide, 0, 0, Inches(0.28), SH, INDIGO)
    # header
    add_text(slide, Inches(0.6), Inches(0.35), Inches(11), Inches(0.4),
             kicker, 13, PURPLE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
             title, 32, DARK, bold=True)
    add_rect(slide, Inches(0.62), Inches(1.55), Inches(2.2), Pt(3), PURPLE)
    # body
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.85),
                                   Inches(12), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, level in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        r = p.add_run()
        if level == 0:
            r.text = "•  " + text
            r.font.size = Pt(18)
            r.font.bold = True
            r.font.color.rgb = DARK
            p.space_after = Pt(6)
            p.space_before = Pt(4)
        else:
            r.text = "–  " + text
            r.font.size = Pt(15)
            r.font.color.rgb = GRAY
            p.space_after = Pt(3)
        r.font.name = FONT
    # footer
    add_text(slide, Inches(0.6), Inches(7.05), Inches(9), Inches(0.3),
             "MindCare — AI Mental Health Chatbot", 10, GRAY)
    add_text(slide, Inches(11.4), Inches(7.05), Inches(1.6), Inches(0.3),
             "Joy Njoroge", 10, GRAY, align=PP_ALIGN.RIGHT)
    return slide


# ============================================================= Slide 1 — Title
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, INDIGO)
add_rect(s, 0, Inches(4.9), SW, Inches(2.6), PURPLE)
add_text(s, Inches(1), Inches(1.6), Inches(11.3), Inches(1.2),
         "🧠  MindCare", 60, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.9), Inches(11.3), Inches(0.8),
         "An AI-Powered Mental Health Chatbot", 28, WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.7),
         "Multilingual (English & Swahili) emotional support, intent classification, "
         "crisis detection & appointment booking", 15, RGBColor(0xE0, 0xDD, 0xFF),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.4), Inches(11.3), Inches(0.5),
         "Joy Njoroge", 20, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
         "Full-Stack AI Project  •  Final Presentation", 14,
         RGBColor(0xEE, 0xEC, 0xFF), align=PP_ALIGN.CENTER)

# ============================================== Slide 2 — Problem Statement
content_slide(
    "Problem Statement", "THE CHALLENGE",
    [("Mental health support in Kenya & East Africa is scarce, costly, and stigmatized", 0),
     ("Very few professionals per capita; long waits and high consultation costs", 1),
     ("Stigma stops many people from seeking help in person", 1),
     ("Most digital mental-health tools are English-only", 0),
     ("This excludes millions of Swahili-first speakers", 1),
     ("People need immediate, private, judgment-free support — any time", 0),
     ("Distress and crisis moments do not wait for office hours", 1),
     ("Crisis situations need instant, safe, responsible guidance", 0),
     ("A wrong or absent response at the wrong moment can be dangerous", 1)],
)

# ======================================================= Slide 3 — Solution
content_slide(
    "Our Solution", "WHAT WE BUILT",
    [("MindCare — a full-stack AI chatbot for emotional support & awareness", 0),
     ("Chats naturally in BOTH English and Swahili (auto-detected)", 0),
     ("Detects the user's language and replies in the same language", 1),
     ("Hybrid intelligence: LLM for natural replies + trained model fallback", 0),
     ("Always responds, even without internet or an API key", 1),
     ("Sentiment-aware, empathetic responses", 0),
     ("Adjusts tone based on detected emotion", 1),
     ("Safety-first crisis detection (English & Swahili)", 0),
     ("Extra services: appointment booking, journaling, therapist directory, M-Pesa donations", 0)],
)

# ============================================== Slide 4 — System Architecture
content_slide(
    "System Architecture", "HOW IT'S STRUCTURED",
    [("Decoupled client–server architecture", 0),
     ("Frontend: React single-page app, deployed on Vercel", 1),
     ("Backend: Flask REST API, deployed on Render (Gunicorn)", 1),
     ("Communication: HTTPS + CORS-scoped JSON/text over REST", 1),
     ("Backend intelligence layers", 0),
     ("ML/NLP engine (intent model, sentiment, entities, language detection)", 1),
     ("LLM providers (Groq / Gemini / Ollama)", 1),
     ("Translation models (MarianMT EN↔SW)", 1),
     ("Data & integrations", 0),
     ("SQLite (appointments)  •  M-Pesa STK Push  •  models pulled from HuggingFace at startup", 1)],
)

# ============================================== Slide 5 — System Design (pipeline)
content_slide(
    "System Design — Chat Pipeline", "REQUEST FLOW",
    [("Every message flows through a layered pipeline:", 0),
     ("1. Detect language — spaCy + Swahili keyword heuristic", 1),
     ("2. If Swahili → translate input to English (MarianMT)", 1),
     ("3. Analyze sentiment (TextBlob) — polarity & mood", 1),
     ("4. Crisis screening — regex (EN + SW) + intent confidence", 1),
     ("5. Generate reply → LLM (primary) or trained model (fallback)", 1),
     ("6. Reply in the user's language (LLM answers in Swahili directly)", 1),
     ("Design principle: layered fallback = the bot always answers", 0),
     ("Crisis gate runs BEFORE generation — safety is never skipped", 1)],
)

# ================================================ Slide 6 — System Modules
content_slide(
    "System Modules", "KEY COMPONENTS",
    [("Chatbot module — intent classification + LLM response generation", 0),
     ("NLP module — sentiment analysis, entity recognition, language detection", 0),
     ("Translation module — English ↔ Swahili (MarianMT), sentence-aware", 0),
     ("Safety module — multilingual crisis detection & safe responses", 0),
     ("Appointment module — book & view counselling sessions (SQLite)", 0),
     ("Donations module — M-Pesa STK Push (Safaricom)", 0),
     ("Frontend modules — chat widget, booking, journal, therapists, resources", 0)],
)

# ================================================== Slide 7 — Tech Stack
content_slide(
    "Project / Tech Stack", "TECHNOLOGIES USED",
    [("Frontend — React 19, Vite 7, Tailwind CSS 4, React Router, lucide-react", 0),
     ("Backend — Python, Flask 3, Flask-CORS, Gunicorn", 0),
     ("Machine Learning — Keras/TensorFlow → TFLite (lightweight inference)", 0),
     ("NLP — NLTK, spaCy, TextBlob", 0),
     ("Translation — HuggingFace Transformers (MarianMT)", 0),
     ("LLM — Groq (gpt-oss-120b), Google Gemini, or Ollama (local)", 0),
     ("Data — SQLite, HuggingFace Hub (model hosting)", 0),
     ("Deployment — Vercel (frontend), Render (backend)", 0)],
)

# =================================================== Slide 8 — Datasets
content_slide(
    "Datasets", "DATA BEHIND THE MODEL",
    [("Intent dataset — hand-authored in intents.json", 0),
     ("155 intents, each with example patterns + curated responses", 1),
     ("458-word vocabulary, 154 output classes after cleaning", 1),
     ("Training data derived via bag-of-words vectorization of patterns", 1),
     ("Translation models — pretrained MarianMT (Rogendo en-sw / sw-en)", 0),
     ("Fine-tuned English↔Swahili neural machine translation", 1),
     ("Model hosting — classifier artifacts on HuggingFace Hub", 0),
     ("Downloaded automatically at startup (kept out of Git)", 1)],
)

# ================================================= Slide 9 — How to Run It
content_slide(
    "How to Run It", "GETTING STARTED",
    [("Backend", 0),
     ("python -m venv venv  →  activate  →  pip install -r requirements.txt", 1),
     ("Set .env (GROQ_API_KEY, etc.)  →  python app.py  →  http://127.0.0.1:5000", 1),
     ("Frontend", 0),
     ("cd frontend  →  npm install  →  npm run dev  →  http://localhost:5173", 1),
     ("Retrain the model (optional)", 0),
     ("python training.py  →  python convert_to_tflite.py", 1),
     ("Deployment", 0),
     ("Backend on Render (gunicorn app:app)  •  Frontend on Vercel", 1)],
)

# ======================================= Slide 10 — Effectiveness & Traction
content_slide(
    "Effectiveness & Traction", "RESULTS SO FAR",
    [("Effectiveness", 0),
     ("Hybrid LLM + trained model → natural yet reliable answers", 1),
     ("Empathy-aware replies tuned by sentiment", 1),
     ("Safety-first crisis handling in two languages", 1),
     ("Works offline via fallback — high availability", 1),
     ("Traction & milestones", 0),
     ("Fully functional full-stack app, deployed and demoable", 1),
     ("Verified bilingual (EN + SW) end-to-end conversations", 1),
     ("Working appointment booking + M-Pesa donations (sandbox)", 1)],
)

# ================================= Slide 11 — Comparison with existing projects
content_slide(
    "Comparison with Existing Solutions", "HOW WE'RE DIFFERENT",
    [("vs. commercial bots (Woebot, Wysa)", 0),
     ("Those are English-only & closed; MindCare is open & Swahili-capable", 1),
     ("vs. simple FAQ / rule-based bots", 0),
     ("MindCare adds LLM reasoning, sentiment & crisis detection", 1),
     ("vs. the earlier version of this project", 0),
     ("Was scikit-learn, English-only, no LLM", 1),
     ("Now: TFLite + LLM + Swahili + crisis safety + booking + payments", 1),
     ("Local context advantage", 0),
     ("Swahili support and M-Pesa donations built for the East African user", 1)],
)

# ============================ Slide 12 — Roadmap, Version Control & Progress
content_slide(
    "Roadmap, Version Control & Status", "WHERE WE ARE & WHAT'S NEXT",
    [("Version control", 0),
     ("Git + GitHub, iterative commits on the main branch", 1),
     ("Current status — up to where we are", 0),
     ("Core chat, bilingual support, booking, donations — all deployed", 1),
     ("Recent work: fixed Swahili detection, model loading & LLM language", 1),
     ("What we want to improve", 0),
     ("User authentication & persistent database (PostgreSQL)", 1),
     ("Conversation memory, email confirmations, admin/therapist portal", 1),
     ("Re-enable semantic search; add more languages; mobile app", 1)],
)

out = "MindCare_Presentation.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
